"""
CONTCAR/POSCAR -> PPT
VESTA CLI a/b/c, POSCAR vs CONTCAR side by side, no stretch
"""
import os, subprocess, time, struct as _struct
from pathlib import Path

from config import VESTA_EXE, LOCAL_BASE
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False

# VESTA_EXE imported from config

AXIS_VIEWS = {
    "a": ((1, 0, 0), (0, 0, 1)),
    "b": ((0, 1, 0), (0, 0, 1)),
    "c": ((0, 0, 1), (0, 1, 0)),
}


def _png_size(path):
    path = Path(path)
    if not path.exists() or path.stat().st_size < 24:
        return None, None
    try:
        with open(str(path), "rb") as f:
            f.read(16)
            w = _struct.unpack(">I", f.read(4))[0]
            h = _struct.unpack(">I", f.read(4))[0]
        return w, h
    except Exception:
        return None, None

def _render_one(struct_path, output_dir, zoom=1.7):
    struct_path = Path(struct_path)
    output_dir = Path(output_dir)
    base = output_dir / "_v_base.vesta"

    subprocess.run(["taskkill", "/IM", "VESTA.exe", "/F"], capture_output=True)

    base_ok = False
    for attempt, wait in enumerate([2.0, 5.0, 8.0]):
        subprocess.Popen(
            [VESTA_EXE, "-open", str(struct_path), "-save", str(base), "-close", ""],
            stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL,
        )
        time.sleep(wait)
        subprocess.run(["taskkill", "/IM", "VESTA.exe", "/F"], capture_output=True)
        if base.exists() and base.stat().st_size > 0:
            base_ok = True
            break
        print("  [retry {}] base.vesta ({}s)".format(attempt + 1, wait))
    if not base_ok:
        return {}

    content = base.read_text(encoding="utf-8")
    content = content.replace("PROJT 0  0.962", "PROJT 0  {:.3f}".format(zoom))
    i = content.find("LORIENT\n")
    e1 = content.find("\n", i) + 1
    e2 = content.find("\n", e1) + 1
    e3 = content.find("\n", e2) + 1
    e4 = content.find("\n", e3) + 1

    result = {}
    for axis, (v1, v2) in AXIS_VIEWS.items():
        r1 = " {:.6f}  {:.6f}  {:.6f}  0.000000  0.000000  0.000000\n".format(*v1)
        r2 = " {:.6f}  {:.6f}  {:.6f}  0.000000  0.000000  0.000000\n".format(*v2)
        cc = content[:e2] + r1 + r2 + content[e4:]
        vf = output_dir / "_v_{}_{}.vesta".format(axis, struct_path.stem)
        vf.write_text(cc, encoding="utf-8")
        pf = output_dir / "_v_{}_{}.png".format(axis, struct_path.stem)

        ok = False
        for attempt, wait in enumerate([2.0, 5.0, 8.0]):
            subprocess.Popen(
                [VESTA_EXE, "-open", str(vf), "-export_img", str(pf), "-close", str(vf)],
                stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL,
            )
            time.sleep(wait)
            subprocess.run(["taskkill", "/IM", "VESTA.exe", "/F"], capture_output=True)
            if pf.exists() and pf.stat().st_size > 0:
                ok = True
                result[axis] = pf
                break
            print("  [retry {}] {} axis ({}s)".format(attempt + 1, axis, wait))

        vf.unlink(missing_ok=True)

    base.unlink(missing_ok=True)
    return result
def add_structure_slides(prs, poscar_dir, zoom=1.7, title=None):
    """a/b axes on one slide (top/bottom), c-axis on separate slide"""
    poscar_dir = Path(poscar_dir)
    poscar_f = poscar_dir / "POSCAR"
    contcar_f = poscar_dir / "CONTCAR"
    if not poscar_f.exists() and not contcar_f.exists():
        print("  [!] no structure files")
        return
    label = title or poscar_dir.name
    out_dir = poscar_dir.parent.parent

    pos = _render_one(poscar_f, out_dir, zoom) if poscar_f.exists() else {}
    ionic_steps = _read_ionic_steps(poscar_dir)
    if ionic_steps < 5:
        print("  [Skip] only {} ionic steps (<5), no meaningful comparison".format(ionic_steps))
        return
    con = _render_one(contcar_f, out_dir, zoom) if contcar_f.exists() else {}

    if not pos and not con:
        print("  [!] render failed")
        return

    def _add_axis_to_slide(slide, axis, p_png, c_png, y_top, half_h):
        """Add one axis row (POSCAR | CONTCAR) at given y position"""
        p_w = p_h = c_w = c_h = None
        if p_png:
            p_w, p_h = _png_size(p_png)
        if c_png:
            c_w, c_h = _png_size(c_png)

        # Calculate max height preserving aspect ratio within half_h
        max_img_h = 0
        for w, h in [(p_w, p_h), (c_w, c_h)]:
            if w and h:
                ratio = h / w
                # Each image gets half the slide width (~6.2 inches)
                if ratio * 5.8 > max_img_h:
                    max_img_h = ratio * 5.8
        img_w = 5.8
        if max_img_h > half_h - 0.3:
            img_w = (half_h - 0.3) / max_img_h * 5.8
            max_img_h = half_h - 0.3

        gap = 0.3
        left_x = (13.333 - (2 * img_w + gap)) / 2

        # Axis label
        axis_label_y = y_top - 0.15
        albl = slide.shapes.add_textbox(Inches(0.3), Inches(axis_label_y), Inches(1.5), Inches(0.15))
        ar = albl.text_frame.paragraphs[0].add_run()
        ar.text = "- {} axis -".format(axis)
        ar.font.size = Pt(10); ar.font.color.rgb = RGBColor(0xBB, 0xBB, 0xBB)
        ar.font.name = "Arial"

        if p_png and p_w and p_h:
            plbl = slide.shapes.add_textbox(Inches(left_x), Inches(y_top - 0.15), Inches(img_w), Inches(0.15))
            prun = plbl.text_frame.paragraphs[0].add_run()
            prun.text = "POSCAR"
            prun.font.size = Pt(8); prun.font.color.rgb = RGBColor(0x90, 0xCA, 0xF9)
            prun.font.name = "Arial"
            plbl.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            slide.shapes.add_picture(str(p_png), Inches(left_x), Inches(y_top), Inches(img_w))
            p_png.unlink()

        if c_png and c_w and c_h:
            rx = left_x + img_w + gap
            clbl = slide.shapes.add_textbox(Inches(rx), Inches(y_top - 0.15), Inches(img_w), Inches(0.15))
            crun = clbl.text_frame.paragraphs[0].add_run()
            crun.text = "CONTCAR"
            crun.font.size = Pt(8); crun.font.color.rgb = RGBColor(0xFF, 0xD7, 0x00)
            crun.font.name = "Arial"
            clbl.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            slide.shapes.add_picture(str(c_png), Inches(rx), Inches(y_top), Inches(img_w))
            c_png.unlink()

    # Slide 1: a+b axes (top/bottom halves)
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])
    slide1.background.fill.solid()
    slide1.background.fill.fore_color.rgb = RGBColor(0x1E, 0x1E, 0x2E)
    tb1 = slide1.shapes.add_textbox(Inches(0.3), Inches(0.05), Inches(12.7), Inches(0.3))
    r1 = tb1.text_frame.paragraphs[0].add_run()
    r1.text = "{} [a/b axis]".format(label)
    r1.font.size = Pt(13); r1.font.bold = True
    r1.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF); r1.font.name = "Arial"

    half_h = 3.4
    _add_axis_to_slide(slide1, "a", pos.get("a"), con.get("a"), 0.4, half_h)
    _add_axis_to_slide(slide1, "b", pos.get("b"), con.get("b"), 0.4 + half_h + 0.1, half_h)

    # Slide 2: c-axis only
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])
    slide2.background.fill.solid()
    slide2.background.fill.fore_color.rgb = RGBColor(0x1E, 0x1E, 0x2E)
    tb2 = slide2.shapes.add_textbox(Inches(0.3), Inches(0.05), Inches(12.7), Inches(0.3))
    r2 = tb2.text_frame.paragraphs[0].add_run()
    r2.text = "{} [c axis]".format(label)
    r2.font.size = Pt(13); r2.font.bold = True
    r2.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF); r2.font.name = "Arial"

    _add_axis_to_slide(slide2, "c", pos.get("c"), con.get("c"), 0.4, 6.8)

    print("  [OK] a+b / c structure slides added")
def add_to_ppt(pptx_path, poscar_dir=None, poscar=None, contcar=None, zoom=1.7):
    if poscar_dir:
        d = Path(poscar_dir)
        prs = Presentation(str(pptx_path)) if Path(pptx_path).exists() else Presentation()
        print("[PPT] {} ({} pages)".format(pptx_path, len(prs.slides)))
        add_structure_slides(prs, d, zoom=zoom)
        prs.save(str(pptx_path))
        print("[OK] saved ({} pages)".format(len(prs.slides)))


if __name__ == "__main__":
    import argparse
    parser = argparse.ArgumentParser()
    parser.add_argument("dir", nargs="?", default=str(LOCAL_BASE / "a_Fe2O3_0701" / "吸附-110_CH_1"))
    parser.add_argument("--pptx", default=str(LOCAL_BASE / "20260623-vasp_report.pptx"))
    parser.add_argument("--zoom", type=float, default=1.7)
    args = parser.parse_args()
    add_to_ppt(args.pptx, poscar_dir=args.dir, zoom=args.zoom)
def _read_ionic_steps(poscar_dir):
    """read .ionic_steps file written by vasp_check"""
    f = Path(poscar_dir) / ".ionic_steps"
    if not f.exists():
        return 99
    try:
        return int(f.read_text(encoding="utf-8").strip())
    except Exception:
        return 99
