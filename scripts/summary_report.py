"""
VASP 汇总报告生成器 (PPT + 结构图)
自动生成 YYYYMMDD-vasp汇总报告.pptx

用法:
  python scripts/summary_report.py --data <json>     # 从 JSON 数据生成
  python scripts/summary_report.py --file <json文件>  # 从 JSON 文件生成
  python scripts/summary_report.py --from-db          # 从项目数据库生成

数据格式:
  [{"proj":"项目名","sub":"子任务名","status":"状态","energy":"...","poscar":"[V]|—","contcar":"[V]|—","task_type":"结构优化/SCF/..."}]
"""
import argparse
import json
import sys
from datetime import datetime
from pathlib import Path
from wcwidth import wcswidth

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from lxml import etree
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False

# 尝试导入 ASE 结构渲染

from struct2ppt import add_structure_slides
from config import LOCAL_BASE, SSH_KEY_DEFAULT
from project_store import load_projects

# LOCAL_BASE imported from config as LOCAL_BASE


# ═══════════════════════════════════════════
#  工具函数
# ═══════════════════════════════════════════

def clean_energy_str(raw):
    if not raw or "TOTEN" not in raw:
        return raw
    try:
        num = raw.split("=")[-1].strip().replace(" eV", "").strip()
        return f"{float(num):.8f}"
    except:
        return raw[:50]


def pad(text, width):
    text = str(text)
    text_width = wcswidth(text)
    if text_width < 0:
        text_width = len(text)
    return text + " " * max(0, width - text_width)


def build_results_from_projects(projects=None, results_file=None):
    """Build results list from project DB, enriched with check results if available"""
    if projects is None:
        projects = load_projects()
    # Load check results for enrichment
    check_data = {}
    if results_file:
        try:
            with open(results_file, "r", encoding="utf-8") as f:
                check_results = json.load(f)
            for r in check_results:
                key = f"{r['proj']}|{r['sub']}"
                check_data[key] = r
        except Exception:
            pass
    """从项目数据库构建 results 列表"""
    if projects is None:
        projects = load_projects()
    results = []
    icons = {
        "Run": "[Run] 运行中", "Completed": "[OK] 已完成",
        "Failed": "[X] 失败", "Error": "[X] 已报错", "Pending": "[~] 待提交", "Stop": "[W] 待接续",
    }
    for proj in projects:
        for sub in proj.get("subs", []):
            # 根据目录前缀推断 task_type
            d = sub.get("dir", "")
            if d.startswith(("opt/", "raw/")):
                tt = "结构优化"
            elif d.startswith(("abs/", "oer/")):
                tt = "结构优化"  # 吸附类也用 IBRION=2 优化
            elif d.startswith("neb/"):
                tt = "NEB"
            elif d.startswith(("static/", "scf/", "dos/", "bands/")):
                tt = "电子结构"
            elif d.startswith("md/"):
                tt = "MD"
            else:
                tt = ""
            results.append({
                "proj": proj["name"], "sub": sub["name"],
                "task_type": tt, "status": icons.get(sub["status"], sub["status"]),
                "energy": "", "poscar": "—", "contcar": "—",
                "dir": d,
            })
    return results


# ═══════════════════════════════════════════
#  结构图渲染 (ASE)
# ═══════════════════════════════════════════

def _find_struct_dirs(projects, results, base_dir):
    """Find dirs needing structure rendering; returns list of (path, title)"""
    struct_dirs = []
    proj_lookup = {p["name"]: p for p in projects}
    for r in results:
        if "[Run]" not in r.get("status", ""):
            continue
        proj = r["proj"]
        sub_name = r.get("sub", "")
        task_type = r.get("task_type", "")
        proj_item = proj_lookup.get(proj)
        if not proj_item:
            continue
        for sub_item in proj_item.get("subs", []):
            if sub_item["name"] != sub_name:
                continue
            candidates = [
                base_dir / proj_item["name"] / sub_name,
                base_dir / proj_item["name"] / sub_item["dir"],
            ]
            for d in candidates:
                if not (d / "POSCAR").exists():
                    continue
                ibrion_str = r.get("ibrion", "")
                if ibrion_str and ibrion_str.lstrip("-").isdigit():
                    ibrion = int(ibrion_str)
                    if ibrion in (2, 3):
                        title = f"{proj}/{sub_item["name"]}/{task_type}"
                        struct_dirs.append((d, title))
                        print(f"  [struct] {title} (IBRION={ibrion})")
                        break
                if task_type and ("结构" in task_type or "optim" in task_type.lower()):
                    title = f"{proj}/{sub_item["name"]}/{task_type}"
                    struct_dirs.append((d, title))
                    print(f"  [struct] {title}")
                    break
            break
    return list(set(struct_dirs))
def write_ppt(results, output_path=None, theme="light", add_struct=True, projects=None):
    """生成 PPT 报告

    Args:
        results: 结果列表
        output_path: 输出路径，默认 {LOCAL_BASE}/YYYYMMDD-vasp汇总报告.pptx
        theme: "light" 或 "dark"
        add_struct: 是否自动添加结构对比图
        projects: 项目数据库（用于查找结构优化目录）
    """
    if not HAS_PPTX:
        print("[!] \u9700\u8981 python-pptx, \u8fd0\u884c: pip install python-pptx")
        return None

    if output_path is None:
        ts = datetime.now().strftime("%Y%m%d")
        output_path = LOCAL_BASE / f"{ts}-vasp\u6c47\u603b\u62a5\u544a.pptx"
    output_path = Path(output_path)

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # ── \u914d\u8272 ──
    if theme == "dark":
        BG, HDR_BG, HDR_FG = RGBColor(0x1E,0x1E,0x2E), RGBColor(0x2D,0x2D,0x44), RGBColor(0xFF,0xFF,0xFF)
        ROW_E, ROW_O = RGBColor(0x28,0x28,0x3C), RGBColor(0x22,0x22,0x38)
        TXT, BDR, PROJ_BG, PROJ_FG = RGBColor(0xE0,0xE0,0xE0), RGBColor(0x3A,0x3A,0x50), RGBColor(0x1A,0x1A,0x2E), RGBColor(0x60,0xA5,0xFA)
    else:
        BG, HDR_BG, HDR_FG = RGBColor(0xFF,0xFF,0xFF), RGBColor(0x1E,0x40,0x73), RGBColor(0xFF,0xFF,0xFF)
        ROW_E, ROW_O = RGBColor(0xF3,0xF4,0xF6), RGBColor(0xFF,0xFF,0xFF)
        TXT, BDR, PROJ_BG, PROJ_FG = RGBColor(0x33,0x33,0x33), RGBColor(0xCC,0xCC,0xCC), RGBColor(0xE8,0xED,0xF3), RGBColor(0x1E,0x40,0x73)

    def st_color(s):
        s = str(s)
        if "[OK]" in s or "\u5df2\u5b8c\u6210" in s: return RGBColor(0x4C,0xAF,0x50)
        if "[X]" in s or "\u5931\u8d25" in s: return RGBColor(0xF4,0x44,0x36)
        if "[Run]" in s or "\u8fd0\u884c" in s: return RGBColor(0x21,0x96,0xF3)
        if "[P]" in s or "\u6392\u961f" in s: return RGBColor(0xFF,0x98,0x00)
        if "[W]" in s or "\u5f85\u63a5\u7eed" in s: return RGBColor(0x9E,0x9E,0x9E)
        return None

    # ── \u5206\u7ec4 ──
    groups = {}
    for r in results:
        p = r.get("proj", "?")
        if p not in groups:
            groups[p] = {"done": [], "run": [], "pend": []}
        st = r.get("status", "")
        if "[OK]" in st or "Completed" in st or "\u5df2\u5b8c\u6210" in st:
            groups[p]["done"].append(r)
        elif "[Run]" in st or "\u8fd0\u884c" in st:
            groups[p]["run"].append(r)
        elif "[P]" in st or "\u6392\u961f" in st:
            groups[p]["pend"].append(r)
        else:
            groups[p]["pend"].append(r)

    cols = ["项目", "子任务", "状态", "类型", "能量(eV)"]
    rows_data = []
    for pname, g in sorted(groups.items()):
        d, run, pend = g["done"], g["run"], g["pend"]
        tc = len(d) + len(run) + len(pend)
        rows_data.append(("proj", [pname, "", f"共{tc}个子任务 | 已完成{len(d)} 运行中{len(run)} 排队{len(pend)}", "", ""]))
        if d:
            ns = "\u3001".join(x["sub"][:10] for x in d[:3])
            if len(d) > 3: ns += f" \u7b49{len(d)}\u4e2a"
            rows_data.append(("sum", ["", ns, "[OK] 已完成", "", ""]))
        for r_item in run:
            rows_data.append(("run", [pname, r_item.get("sub",""), "[Run] 运行中", r_item.get("task_type","-"), r_item.get("energy","-")]))
        if pend:
            run_pend = [r_item for r_item in pend if "[P]" in r_item.get("status","") or "\u6392\u961f" in r_item.get("status","")]
            waiting = [r_item for r_item in pend if r_item not in run_pend]
            if run_pend:
                ns = "\u3001".join(x["sub"][:10] for x in run_pend[:3])
                if len(run_pend) > 3: ns += f" \u7b49{len(run_pend)}\u4e2a"
                rows_data.append(("sum", ["", ns, "[P] 排队中", "", ""]))
            if waiting:
                ns = "\u3001".join(x["sub"][:10] for x in waiting[:3])
                if len(waiting) > 3: ns += f" \u7b49{len(waiting)}\u4e2a"
                rows_data.append(("sum", ["", ns, "[W] 待接续", "", ""]))

    # ── \u5c01\u9762 ──
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    sl.background.fill.solid(); sl.background.fill.fore_color.rgb = BG
    tb = sl.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11), Inches(1.5))
    r0 = tb.text_frame.paragraphs[0].add_run()
    r0.text = "VASP \u9879\u76ee\u72b6\u6001\u6c47\u603b\u62a5\u544a"
    r0.font.size = Pt(36); r0.font.bold = True; r0.font.color.rgb = PROJ_FG; r0.font.name = "\u5fae\u8f6f\u96c5\u9ed1"
    tb.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    tb2 = sl.shapes.add_textbox(Inches(1), Inches(4), Inches(11), Inches(0.6))
    r1 = tb2.text_frame.paragraphs[0].add_run()
    r1.text = f"\u751f\u6210\u65f6\u95f4: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}    \u5171{len(results)}\u4e2a\u5b50\u4efb\u52a1"
    r1.font.size = Pt(14); r1.font.color.rgb = TXT; r1.font.name = "\u5fae\u8f6f\u96c5\u9ed1"
    tb2.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # ── \u5206\u9879\u76ee\u9875 ──
    proj_pages = []
    for kind, cells in rows_data:
        if kind == "proj":
            proj_pages.append([])
        if proj_pages:
            proj_pages[-1].append((kind, cells))

    cw = [Inches(1.5), Inches(4.5), Inches(2.0), Inches(1.5), Inches(3.0)]

    def bd(cell):
        tcPr = cell._tc.get_or_add_tcPr()
        ns = 'http://schemas.openxmlformats.org/drawingml/2006/main'
        for s in ['lnL','lnR','lnT','lnB']:
            ln = tcPr.find(etree.QName(ns, s))
            if ln is None: ln = etree.SubElement(tcPr, etree.QName(ns, s))
            ln.set('w', '6350')
            sf = ln.find(etree.QName(ns, 'solidFill'))
            if sf is None: sf = etree.SubElement(ln, etree.QName(ns, 'solidFill'))
            sc = sf.find(etree.QName(ns, 'srgbClr'))
            if sc is None: sc = etree.SubElement(sf, etree.QName(ns, 'srgbClr'))
            sc.set('val', f"{BDR[0]:02X}{BDR[1]:02X}{BDR[2]:02X}")

    def fc(cell, txt, bold=False, sz=9, align=PP_ALIGN.LEFT, color=None, bg=None):
        cell.text = ""
        p = cell.text_frame.paragraphs[0]; r = p.add_run()
        r.text = str(txt); r.font.size = Pt(sz); r.font.bold = bold
        r.font.name = "\u5fae\u8f6f\u96c5\u9ed1"; p.alignment = align
        if color: r.font.color.rgb = color
        cell.vertical_anchor = 1
        if bg: cell.fill.solid(); cell.fill.fore_color.rgb = bg

    for page in proj_pages:
        sl2 = prs.slides.add_slide(prs.slide_layouts[6])
        sl2.background.fill.solid(); sl2.background.fill.fore_color.rgb = BG
        n_pr = len(page) + 1
        rh = 0.40 if n_pr <= 8 else (0.34 if n_pr <= 15 else 0.28)
        tbl = sl2.shapes.add_table(n_pr, len(cols), Inches(0.3), Inches(0.5), Inches(12.5), Inches(min(n_pr*rh+0.5, 6.5))).table
        for i, w in enumerate(cw): tbl.columns[i].width = w

        for j, h in enumerate(cols):
            c = tbl.cell(0, j)
            fc(c, h, bold=True, sz=10, align=PP_ALIGN.CENTER, color=HDR_FG, bg=HDR_BG)
            bd(c)

        for ri, (kind, cells) in enumerate(page, 1):
            bg_r = PROJ_BG if kind == "proj" else (ROW_E if ri % 2 == 0 else ROW_O)
            for j, val in enumerate(cells):
                c = tbl.cell(ri, j)
                sc = st_color(val) if j == 1 else None
                if kind == "proj":
                    fc(c, val, bold=True, sz=10, align=PP_ALIGN.LEFT, color=PROJ_FG, bg=bg_r)
                elif kind == "sum":
                    fc(c, val, bold=(j==1), sz=9, align=PP_ALIGN.LEFT, color=sc or TXT, bg=bg_r)
                else:
                    a = PP_ALIGN.CENTER if j in (0,4) else PP_ALIGN.LEFT
                    fc(c, val, sz=9, align=a, color=sc or TXT, bg=bg_r)
                bd(c)
                if kind == "proj" and j == 0:
                    try: tbl.cell(ri, 0).merge(tbl.cell(ri, len(cols)-1))
                    except: pass
                if kind == "proj" and j > 0:
                    continue

    # ── \u7ed3\u6784\u56fe (规则: 运行中的项目下有结构优化子任务) ──
    _ssh_client = None
    if add_struct:
        if projects is None:
            projects = load_projects()
        try:
            sys.path.insert(0, str(Path(__file__).resolve().parent))
            from ssh_connect import parse_host
            u, h = parse_host("mdye@hpc.xmu.edu.cn")
            _ssh_client = create_ssh_client(h, 22, u, str(SSH_KEY_DEFAULT), None, 10)
        except Exception:
            pass
            projects = load_projects()
        struct_dirs = _find_struct_dirs(projects, results, LOCAL_BASE)
        for sd, sd_title in struct_dirs:
            add_structure_slides(prs, sd, title=sd_title)

    prs.save(str(output_path))
    print(f"[OK] PPT: {output_path} ({len(prs.slides)} \u9875)")
    return output_path


def write_summary(results, output_dir=None):
    """\u751f\u6210\u6c47\u603b\u62a5\u544a (\u7ec8\u7aef\u6253\u5370 + PPT)"""
    if output_dir is None:
        output_dir = LOCAL_BASE
    output_dir = Path(output_dir)
    output_dir.mkdir(parents=True, exist_ok=True)

    # \u7ec8\u7aef\u6253\u5370
    print(f"\n{'='*60}")
    print(f" VASP \u4f5c\u4e1a\u68c0\u67e5\u6c47\u603b (\u5171{len(results)}\u4e2a\u5b50\u4efb\u52a1)")
    print(f"{'='*60}")
    header = (pad("\u9879\u76ee", 14) + " " + pad("\u5b50\u4efb\u52a1", 16) + " "
              + pad("\u7c7b\u578b", 10) + " " + pad("\u72b6\u6001", 16) + " "
              + pad("\u80fd\u91cf(eV)", 40) + " " + pad("POSCAR", 8) + " " + "CONTCAR")
    sep = "-" * 112
    print(header); print(sep)
    for r in results:
        icon = f"[X] {r['status']}" if r["status"] not in ("\u5df2\u6536\u655b", "\u672a\u6536\u655b") else ("[OK] \u5df2\u6536\u655b" if r["status"] == "\u5df2\u6536\u655b" else "[!] \u672a\u6536\u655b")
        raw = r["energy"]
        energy_str = clean_energy_str(raw) if raw and "TOTEN" in raw else (raw[:38] if raw else "\u2014")
        line = (pad(r['proj'], 14) + " " + pad(r['sub'], 16) + " "
                + pad(r.get("task_type",""), 10) + " " + pad(icon, 16) + " "
                + pad(energy_str, 40) + " " + pad(r['poscar'], 8) + " " + r['contcar'])
        print(line)
    print(sep)
    print(f"\u5171\u68c0\u67e5 {len(results)} \u4e2a\u5b50\u4efb\u52a1")

    # \u751f\u6210 PPT
    ts = datetime.now().strftime("%Y%m%d")
    ppt_path = write_ppt(results, output_dir / f"{ts}-vasp\u6c47\u603b\u62a5\u544a.pptx")
    return ppt_path


def main():
    parser = argparse.ArgumentParser(description="VASP \u6c47\u603b\u62a5\u544a\u751f\u6210\u5668")
    parser.add_argument("--data", default=None, help="JSON \u5b57\u7b26\u4e32")
    parser.add_argument("--file", default=None, help="JSON \u6587\u4ef6\u8def\u5f84")
    parser.add_argument("--from-db", action="store_true", help="\u4ece\u9879\u76ee\u6570\u636e\u5e93\u751f\u6210")
    parser.add_argument("--from-results", default=None, help="\u4ece\u68c0\u67e5\u7ed3\u679cJSON\u751f\u6210\uff08\u4f18\u5148\u4e8e--from-db\uff09")
    parser.add_argument("-o", "--output", default=None, help="\u8f93\u51fa\u8def\u5f84")
    parser.add_argument("--theme", choices=["dark", "light"], default="light", help="PPT \u4e3b\u9898")
    parser.add_argument("--no-struct", action="store_true", help="\u4e0d\u6dfb\u52a0\u7ed3\u6784\u56fe")
    args = parser.parse_args()

    results = None
    projects = None
    if args.data:
        results = json.loads(args.data)
    elif args.file:
        with open(args.file, "r", encoding="utf-8") as f:
            results = json.load(f)
    elif args.from_db:
        projects = load_projects()
        results = build_results_from_projects(projects)
    else:
        projects = load_projects()
        results = build_results_from_projects(projects)
        if not results:
            print("[!] \u65e0\u6570\u636e"); sys.exit(1)

    if not results:
        print("[!] \u65e0\u6570\u636e"); sys.exit(1)

    print(f"[OK] \u8bfb\u53d6 {len(results)} \u6761\u8bb0\u5f55")
    write_ppt(results, args.output, args.theme, add_struct=not args.no_struct, projects=projects)


if __name__ == "__main__":
    main()